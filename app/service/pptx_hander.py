import logging
from fastapi import HTTPException
import json
from pptx import Presentation
from app.service.knowledgebase_handler import get_knowledge_base
from app.utils.langchain_utils import process_user_question

async def generate_presentation():
    context = get_knowledge_base()
    if not context:
        raise HTTPException(status_code=400, detail="Please upload a PDF file first.")

    
    # Create the prompt by combining the requirement for PPTX and PDF context
    prompt = f"Write a powerpoint presentation about this document in JSON. The JSON should contain 'slides' and each slide should contain 'title' and 'content'. Write the answer only in JSON code without any other thing."

    try:
        response_str = process_user_question(prompt)
        json_start_index = response_str.index('"slides":')
        json_string = "{"+response_str[json_start_index:]# Extract the JSON string from the response
        json_pptx = json.loads(json_string)
        output_file = process_presentation(json_pptx)
        return output_file
    except Exception as e:
        logging.error(f"An error occurred during generate_presentation(): {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

def process_presentation(json_data):
    try:
        validate_json_data_structure(json_data)
        presentation = Presentation()
        slides = json_data.get("slides", []) if isinstance(json_data, dict) else json_data

        for slide_data in slides:
            if isinstance(slide_data, dict):
                title = slide_data.get("title", "")
                content = slide_data.get("content", "")

                slide_layout = presentation.slide_layouts[1]

                slide = presentation.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                slide.placeholders[1].text = content

        output_file = "presentation.pptx"
        presentation.save(output_file)

        return output_file
    except Exception as e:
        logging.error(f"An error occurred during process_presentation(): {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))

def validate_json_data_structure(json_data):
    if not isinstance(json_data, (dict, list)):
        raise ValueError("Invalid JSON data format")

    if isinstance(json_data, dict):
        slides = json_data.get("slides")
        if slides is None:
            raise ValueError("Missing 'slides' field in slide data")
    else:
        slides = json_data

    for slide_data in slides:
        if not isinstance(slide_data, dict):
            raise ValueError("Invalid JSON data format")
        if "title" not in slide_data:
            raise ValueError("Missing 'title' field in slide data")
        if not isinstance(slide_data["title"], str):
            raise ValueError("'title' field must be a string")
        if "content" not in slide_data:
            raise ValueError("Missing 'content' field in slide data")
        if not isinstance(slide_data["content"], str):
            raise ValueError("'content' field must be a string")